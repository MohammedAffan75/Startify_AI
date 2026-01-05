"""Assembler: combine outputs from research/generation into final package."""
import os
import json
import zipfile
from typing import Dict, Any
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_CENTER, TA_LEFT


def assemble_package(idea_struct: dict, research_results: dict, branding_content: dict, output_dir: str) -> str:
    """Assemble complete package with PPTX, PDF, and JSON assets.
    
    Args:
        idea_struct: Parsed idea structure
        research_results: Research findings
        branding_content: Generated branding and content
        output_dir: Output directory path (e.g., 'outputs/<job_id>')
        
    Returns:
        Path to the main output file (pitch_deck.pptx)
    """
    # Create output directory
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    
    # Generate pitch deck PPTX
    pptx_path = os.path.join(output_dir, "pitch_deck.pptx")
    _create_pitch_deck(idea_struct, research_results, branding_content, pptx_path)
    
    # Generate summary PDF
    pdf_path = os.path.join(output_dir, "summary.pdf")
    _create_summary_pdf(idea_struct, research_results, branding_content, pdf_path)
    
    # Save assets JSON
    assets_path = os.path.join(output_dir, "assets.json")
    _save_assets_json(branding_content, assets_path)
    
    return pptx_path


def _create_pitch_deck(idea_struct: dict, research_results: dict, branding_content: dict, output_path: str):
    """Create pitch deck PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    pitch_sections = branding_content.get("pitch_sections", {})
    brand_names = branding_content.get("brand_names", [])
    slogans = branding_content.get("slogans", [])
    
    # Slide 1: Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = brand_names[0] if brand_names else "Startup Idea"
    subtitle.text = slogans[0] if slogans else "Innovative Solution"
    
    # Slide 2: Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "The Problem"
    tf = content.text_frame
    tf.text = pitch_sections.get("problem", "Market problem statement")
    
    # Slide 3: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Our Solution"
    tf = content.text_frame
    tf.text = pitch_sections.get("solution", "Solution description")
    
    # Slide 4: Market & Trends
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Market Opportunity"
    tf = content.text_frame
    
    # Add market size
    p = tf.paragraphs[0]
    p.text = pitch_sections.get("market_size_estimate", "Market analysis")
    
    # Add competitors
    competitors = research_results.get("competitors", [])
    if competitors:
        p = tf.add_paragraph()
        p.text = "\nKey Competitors:"
        p.level = 0
        for comp in competitors[:3]:
            p = tf.add_paragraph()
            p.text = f"• {comp.get('name', 'Unknown')}"
            p.level = 1
    
    # Slide 5: Business Model
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Business Model"
    tf = content.text_frame
    tf.text = pitch_sections.get("business_model", "Revenue model")
    
    # Slide 6: Financials Placeholder
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Financial Projections"
    tf = content.text_frame
    tf.text = """Year 1: $100K ARR (projected)
Year 2: $500K ARR (projected)
Year 3: $2M ARR (projected)

Key Metrics:
• Customer Acquisition Cost: TBD
• Lifetime Value: TBD
• Gross Margin: 70%+"""
    
    # Slide 7: Team & Ask
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Team & The Ask"
    tf = content.text_frame
    
    p = tf.paragraphs[0]
    p.text = pitch_sections.get("team_reqs", "Team requirements")
    
    p = tf.add_paragraph()
    p.text = "\n\nThe Ask:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Seeking $500K seed funding for:"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Product development & MVP launch"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Initial marketing & user acquisition"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Team expansion"
    p.level = 2
    
    # Save presentation
    prs.save(output_path)


def _create_summary_pdf(idea_struct: dict, research_results: dict, branding_content: dict, output_path: str):
    """Create one-page summary PDF."""
    doc = SimpleDocTemplate(output_path, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor='#2C3E50',
        spaceAfter=30,
        alignment=TA_CENTER
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=14,
        textColor='#34495E',
        spaceAfter=12,
        spaceBefore=12
    )
    
    # Title
    brand_names = branding_content.get("brand_names", [])
    title = brand_names[0] if brand_names else "Startup Idea Summary"
    story.append(Paragraph(title, title_style))
    story.append(Spacer(1, 0.2*inch))
    
    # Brand Names Section
    story.append(Paragraph("Brand Name Options", heading_style))
    brand_text = "<br/>".join([f"• {name}" for name in brand_names[:5]])
    story.append(Paragraph(brand_text, styles['Normal']))
    story.append(Spacer(1, 0.2*inch))
    
    # Slogans
    story.append(Paragraph("Tagline Options", heading_style))
    slogans = branding_content.get("slogans", [])
    slogan_text = "<br/>".join([f"• {slogan}" for slogan in slogans[:3]])
    story.append(Paragraph(slogan_text, styles['Normal']))
    story.append(Spacer(1, 0.2*inch))
    
    # Top 3 Insights
    story.append(Paragraph("Key Insights", heading_style))
    
    opportunities = research_results.get("key_opportunities", [])
    insights_text = "<br/>".join([f"• {opp}" for opp in opportunities[:3]])
    story.append(Paragraph(insights_text, styles['Normal']))
    story.append(Spacer(1, 0.2*inch))
    
    # Market Trends
    story.append(Paragraph("Market Trends", heading_style))
    trends = research_results.get("trends", {})
    if trends:
        trend_text = "<br/>".join([f"• {k}: {v:.1f}/100" for k, v in list(trends.items())[:3]])
        story.append(Paragraph(trend_text, styles['Normal']))
    else:
        story.append(Paragraph("Trend data not available", styles['Normal']))
    
    # Build PDF
    doc.build(story)


def _save_assets_json(branding_content: dict, output_path: str):
    """Save branding assets as JSON."""
    assets = {
        "brand_names": branding_content.get("brand_names", []),
        "slogans": branding_content.get("slogans", []),
        "logo_prompts": branding_content.get("logo_prompts", []),
        "ad_copies": branding_content.get("ad_copies", []),
        "pitch_sections": branding_content.get("pitch_sections", {})
    }
    
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(assets, f, indent=2, ensure_ascii=False)


def create_downloadable_zip(job_id: str, base_output_dir: str = "outputs") -> str:
    """Create a zip file of the job output directory.
    
    Args:
        job_id: Job identifier
        base_output_dir: Base directory where outputs are stored
        
    Returns:
        Path to the created zip file
    """
    output_dir = os.path.join(base_output_dir, job_id)
    zip_path = os.path.join(base_output_dir, f"{job_id}.zip")
    
    if not os.path.exists(output_dir):
        raise FileNotFoundError(f"Output directory not found: {output_dir}")
    
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
        for root, dirs, files in os.walk(output_dir):
            for file in files:
                file_path = os.path.join(root, file)
                arcname = os.path.relpath(file_path, output_dir)
                zipf.write(file_path, arcname)
    
    return zip_path
